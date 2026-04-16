import base64
from collections.abc import Iterable
import hashlib
from logging import Logger
import os
from pathlib import Path
import requests
import shutil
import subprocess
import tempfile
from typing import Optional

from docx2python import docx2python
import pymupdf
from pptx import Presentation
from PIL import Image
from pydantic_settings import BaseSettings

import fitz
import io
import re

from collections import OrderedDict
import threading

class FileWorker:
    def __init__(
            self, 
            logger: Logger, 
            file, 
            format: str, 
            settings: BaseSettings,
            diarization_params: dict
        ):
        self.settings = settings
        self.logger = logger
        self.file = file
        self.diarization_params = diarization_params
        self.format = format.lower()
        self._document_cache = OrderedDict()
        self._cache_lock = threading.Lock()


    def text_extractor(self) -> str:
        self.logger.info(f"Processing file with format {self.format}")
        
        file_hash = None
        try:
            file_hash = self._compute_file_hash(self.file)
            if (cached := self._get_cached_text(file_hash, self.format)) is not None:
                self.logger.info(f"Cache HIT for {self.file} ({self.format})")
                return cached
        except Exception as e:
            self.logger.warning(f"Failed to hash file, skipping cache: {e}")

        try:
            if self.format in [".jpg", ".jpeg", ".png", ".bmp", ".gif", ".tiff"]:
                text = self._extract_text_from_image()
            elif self.format == ".pdf":
                text = self._extract_text_from_pdf()
            elif self.format in [".docx", ".doc"]:
                text = self._extract_text_from_word()
            elif self.format == ".pptx":
                text = self._extract_text_from_pptx()
            elif self.format == ".emf":
                text = self._extract_text_from_emf()
            elif self.format in [
                ".mp3", ".wav", ".flac", ".ogg", ".webm", ".mov", ".mkv", ".avi", ".mp4"
            ]:
                text = self._extract_text_from_media()
            else:
                error_msg = f"Unsupported format: {self.format}. No handler implemented in FileWorker."
                self.logger.error(error_msg)
                raise ValueError(error_msg)

            if file_hash is not None and text.strip():
                self._set_cached_text(file_hash, self.format, text)
                self.logger.debug(f"Cached result for {file_hash}:{self.format}")

            self.logger.info(f"Successfully extracted text, length: {len(text)} characters")
            return text

        except Exception:
            self.logger.exception(f"Unhandled exception during {self.format} extraction")
            raise

    def _compute_file_hash(self, file_path: str | Path, chunk_size: int = 8192) -> str:
        """
        Compute SHA256 hash of file content in chunks to handle large files efficiently.
        
        :param file_path: Path to the file
        :param chunk_size: Read buffer size in bytes
        :return: Hexadecimal SHA256 hash string
        :raises FileNotFoundError: If file does not exist
        :raises IOError: On read errors
        """
        sha256 = hashlib.sha256()
        with open(file_path, 'rb') as f:
            while chunk := f.read(chunk_size):
                sha256.update(chunk)
        return sha256.hexdigest()
    
    def _get_cached_text(self, file_hash: str, file_format: str) -> Optional[str]:
        """
        Retrieve cached extraction result for given file hash and format.
        
        :param file_hash: SHA256 hash of file content
        :param file_format: File extension (e.g., '.pdf')
        :return: Cached text if exists, None otherwise
        """
        cache_key = f"{file_hash}:{file_format}"
        with self._cache_lock:
            if cache_key in self._document_cache:
                self._document_cache.move_to_end(cache_key)
                return self._document_cache[cache_key]
        return None
    
    def _set_cached_text(self, file_hash: str, file_format: str, text: str) -> None:
        """
        Store extraction result in cache with LRU eviction policy.
        
        :param file_hash: SHA256 hash of file content
        :param file_format: File extension
        :param text: Extracted text to cache
        """
        cache_key = f"{file_hash}:{file_format}"
        with self._cache_lock:
            self._document_cache[cache_key] = text
            self._document_cache.move_to_end(cache_key)
            # Evict oldest item if over capacity
            if len(self._document_cache) > self.settings.CACHE_MAXSIZE:
                self._document_cache.popitem(last=False)

    def _extract_text_from_pdf(self) -> str:
        """
        Extract text from PDF via page-by-page rasterization and Ollama vision analysis.
        :return: Concatenated text with page markers (Russian labels preserved per business logic)
        :raises Exception: Propagates critical PDF processing failures
        """
        full_text = []
        with pymupdf.open(self.file) as doc:
            pages = len(doc)
            if pages > self.settings.PDF_PAGES_LIMIT:
                raise ValueError(f"Page limit exceeded: {pages} > {self.settings.PDF_PAGES_LIMIT}")
            self.logger.info(f"Document contains {pages} pages")

            for page_num in range(pages):
                tmp_img_path = None
                try:
                    page = doc.load_page(page_num)
                    pix = page.get_pixmap(dpi=self.settings.DPI)
                    
                    # Save page as temporary PNG
                    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp_img:
                        tmp_img_path = tmp_img.name
                        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                        img.save(tmp_img_path, format="PNG")
                    
                    # Analyze page image with vision model
                    page_text = self._send_image_to_ollama(tmp_img_path)
                    
                    if page_text and not page_text.startswith("["):
                        full_text.append(f"--- Page {page_num + 1} ---\n{page_text}")
                        self.logger.info(f"Page {page_num + 1}: extracted {len(page_text)} chars")
                    else:
                        full_text.append(f"--- Page {page_num + 1} [error] ---\n{page_text}")
                        self.logger.warning(f"Page {page_num + 1} analysis warning: {page_text}")
                
                except Exception as e:
                    error_msg = f"[Page processing error: {str(e)[:100]}]"
                    full_text.append(f"--- Page {page_num + 1} [crash] ---\n{error_msg}")
                    self.logger.error(f"Page {page_num + 1} processing crashed", exc_info=True)
                finally:
                    # Critical fix: guaranteed temp image cleanup on all paths
                    if tmp_img_path and os.path.exists(tmp_img_path):
                        try:
                            os.unlink(tmp_img_path)
                        except OSError as e:
                            self.logger.warning(f"Failed to delete temp image for page {page_num + 1}: {e}")
        
        # Post-processing: normalize line breaks
        result = "\n\n".join(full_text)
        result = re.sub(r'\n{3,}', '\n\n', result)
        total_chars = len(result)
        
        self.logger.info(f"PDF processing complete: {pages} pages, {total_chars} chars")
        if total_chars < 50:
            return f"[PDF poorly recognized, only {total_chars} characters]\n{result}"
        return result

    def _convert_emf_to_png(self, emf_path: str) -> str | None:
        """
        Convert EMF vector graphic to PNG using LibreOffice headless mode.
        
        :param emf_path: Path to source EMF file
        :return: Path to converted PNG file, or None on failure
        :note: Requires LibreOffice installed and available in system PATH
        """
        emf_path = Path(emf_path)
        png_path = emf_path.with_suffix(".png")
        
        # Critical fix: verify dependency availability before execution
        if not shutil.which("libreoffice"):
            self.logger.error(
                "LibreOffice not found in PATH. EMF conversion requires LibreOffice installation."
            )
            return None
        
        try:
            result = subprocess.run(
                [
                    "libreoffice",
                    "--headless",
                    "--convert-to", "png",
                    "--outdir", str(emf_path.parent),
                    str(emf_path)
                ],
                capture_output=True,
                text=True,
                timeout=60
            )
            
            if result.returncode == 0 and png_path.exists():
                return str(png_path)
            else:
                self.logger.error(
                    f"LibreOffice conversion failed (code {result.returncode}). "
                    f"Stderr: {result.stderr[:200]}"
                )
                return None
                
        except subprocess.TimeoutExpired:
            self.logger.error("LibreOffice conversion timed out after 60 seconds")
            return None
        except Exception as e:
            self.logger.error(f"Unexpected EMF conversion error: {e}", exc_info=True)
            return None

    def _extract_text_from_image(self) -> str:
        """Extract text from standalone image file using Ollama."""
        description = self._send_image_to_ollama(self.file)
        if description and not description.startswith("["):
            return f"📸 **IMAGE ANALYSIS:**\n\n{description}"
        return description

    def _extract_text_from_word(self) -> str:
        """Extract text from Word documents (.doc and .docx)"""
        try:
            if self.format == ".doc":
                # Convert .doc to .docx
                docx_path = self._convert_doc_to_docx(self.file)
                if not docx_path:
                    return "[Error: failed to convert .doc to .docx]"
                
                try:
                    # Use existing .docx pipeline
                    with docx2python(docx_path) as doc_result:
                        all_parts = [doc_result.body, doc_result.header, doc_result.footer]
                    return self._word_to_text(all_parts)
                finally:
                    # Remove temporary file
                    if os.path.exists(docx_path):
                        os.unlink(docx_path)
            else:
                # Direct .docx processing
                with docx2python(self.file) as doc_result:
                    all_parts = [doc_result.body, doc_result.header, doc_result.footer]
                return self._word_to_text(all_parts)
                
        except Exception as e:
            self.logger.error(f"Error processing Word document: {e}")
            return ""

    def _convert_doc_to_docx(self, doc_path: str | Path) -> str | None:
        """
        Convert .doc to .docx using LibreOffice/unoconv

        :param doc_path: Path to source .doc file
        :return: Path to converted .docx file or None on error
        """
        doc_path = Path(doc_path)

        # Check LibreOffice availability
        if not shutil.which("libreoffice"):
            self.logger.error("LibreOffice not found. Required for .doc conversion.")
            return None

        # Create temporary file for result
        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
            docx_path = tmp.name

        try:
            # Use LibreOffice for conversion
            result = subprocess.run(
                [
                    "libreoffice",
                    "--headless",
                    "--convert-to", "docx",
                    "--outdir", str(doc_path.parent),
                    str(doc_path)
                ],
                capture_output=True,
                text=True,
                timeout=60
            )

            if result.returncode == 0:
                # LibreOffice creates file with same name but .docx extension
                auto_converted = doc_path.with_suffix(".docx")
                if auto_converted.exists():
                    # Move to our temporary file
                    shutil.move(str(auto_converted), docx_path)
                    return docx_path
                else:
                    # Search for file in output directory
                    for file in doc_path.parent.glob(f"{doc_path.stem}*.docx"):
                        shutil.move(str(file), docx_path)
                        return docx_path

            self.logger.error(f"LibreOffice conversion failed: {result.stderr[:200]}")
            return None

        except subprocess.TimeoutExpired:
            self.logger.error("LibreOffice conversion timed out")
            return None
        except Exception as e:
            self.logger.error(f"Unexpected conversion error: {e}", exc_info=True)
            return None
        finally:
            # If temporary file was not used, delete it
            if os.path.exists(docx_path):
                try:
                    os.unlink(docx_path)
                except Exception as e:
                    self.logger.warning(f"Failed to cleanup temp docx file: {e}")

    def _send_image_to_ollama(self, image_path: str | Path) -> str:
        """
        Sends an image file to Ollama vision model and returns extracted text.
        
        :param image_path: Path to image file (PNG, JPG, etc.)
        :return: Extracted description or error message
        """
        try:
            # Read and encode image
            with open(image_path, "rb") as f:
                img_bytes = f.read()
            img_base64 = base64.b64encode(img_bytes).decode('utf-8')

            payload = {
                "model": self.settings.OLLAMA_VISION_MODEL,
                "prompt": self.settings.IMAGE_PROMPT,
                "images": [img_base64],
                "stream": False,
                "options": {"temperature": 0.1},
            }

            response = requests.post(
                self.settings.OLLAMA_URL,
                json=payload,
                timeout=120
            )

            if response.status_code == 200:
                result = response.json()
                description = result.get("response", "").strip()
                return description if description else "[Image processed, but no description received]"
            else:
                error_detail = response.json().get("error", "") if response.headers.get("content-type", "").startswith("application/json") else response.text[:100]
                return f"[Ollama error {response.status_code}: {error_detail}]"

        except requests.exceptions.Timeout:
            return "[Timeout when connecting to Ollama]"
        except FileNotFoundError:
            return "[Image file not found]"
        except Exception as e:
            return f"[Error sending to Ollama: {str(e)}]"

    def _safe_decode(self, s: str) -> str:
        """
        Decode Word metadata if necessary.
        Only decodes if the result contains more Cyrillic characters than the original string.
        """
        if not isinstance(s, str):
            return s
        try:
            decoded = s.encode('latin1').decode('cp1251')
            cyrillic = lambda t: sum('\u0400' <= c <= '\u04FF' for c in t)
            if cyrillic(decoded) > cyrillic(s):
                return decoded
            return s
        except (UnicodeEncodeError, UnicodeDecodeError):
            return s
        
    def _extract_text_from_pptx(self) -> str:
        try:
            prs = Presentation(self.file)
            text_lines = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())
                
                if slide_content:
                    text_lines.append(f"--- Slide {slide_num} ---")
                    text_lines.extend(slide_content)
            
            return "\n".join(text_lines)
        except Exception as e:
            self.logger.error(f"Error processing PowerPoint: {e}")
            return ""

    def _word_to_text(self, all_parts: list[list[list[list[list[str]]]]]) -> str:
        """
        Extract text from Word document structure
        """
        text_items = []

        def extract_text_recursively(data: str | Iterable) -> None:
            if isinstance(data, str):
                if data and data.strip():
                    # Decode if necessary
                    decoded = self._safe_decode(data.strip())
                    text_items.append(decoded)
            elif isinstance(data, Iterable):
                for item in data:
                    extract_text_recursively(item)

        for part in all_parts:
            extract_text_recursively(part)

        # Filter empty strings and join
        return '\n'.join(filter(None, text_items))
    
    def _extract_text_from_emf(self) -> str:
        """Convert EMF to PNG and extract text via Ollama."""
        png_path = self._convert_emf_to_png(self.file)
        if not png_path:
            return "[Error: failed to convert EMF to PNG]"

        try:
            description = self._send_image_to_ollama(png_path)
            if description and not description.startswith("["):
                return f"🖼️ **VECTOR IMAGE ANALYSIS (EMF):**\n\n{description}"
            return description
        finally:
            try:
                Path(png_path).unlink(missing_ok=True)
            except Exception as e:
                self.logger.warning(f"Failed to remove temp PNG after EMF conversion: {e}")

    def _extract_text_from_media(self) -> str:
        try:
            with open(self.file, "rb") as fh:
                file_bytes = fh.read()
        except Exception as e:
            return str(e)

        try:
            files = {"file": (self.file.name, file_bytes, "application/octet-stream")}

            response = requests.post(
                self.settings.WHISPER_API_URL,
                params=self.diarization_params,
                files=files,
                timeout=self.settings.WHISPER_TIMEOUT,
            )
        except Exception as e:
            self.logger.error(f"Error during transcription request: {e}") 
            return str(e)
        if response.status_code == 200:
            result = response.json()
            transcription = result.get("result", str(result)) 
        else:
            transcription = f"Transcription failed with status {response.status_code}"

        return transcription
